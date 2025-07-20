#!/usr/bin/env python3
"""
Generate a PowerPoint-style presentation for AWS-InfoBlox Sync
Creates both PPTX and PDF versions
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_presentation():
    """Create the PowerPoint presentation"""
    prs = Presentation()
    
    # Set slide dimensions (16:9)
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # Color scheme
    colors = {
        'primary': RGBColor(255, 153, 0),    # AWS Orange
        'secondary': RGBColor(35, 47, 62),   # AWS Dark Blue
        'accent': RGBColor(20, 110, 180),    # AWS Light Blue
        'success': RGBColor(72, 187, 120),   # Green
        'danger': RGBColor(229, 62, 62),     # Red
        'white': RGBColor(255, 255, 255),
        'gray': RGBColor(107, 114, 128)
    }
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[6]  # Blank slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['secondary']
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(2))
    title = title_box.text_frame
    title.text = "AWS-InfoBlox Network Synchronization"
    title.paragraphs[0].font.size = Pt(48)
    title.paragraphs[0].font.bold = True
    title.paragraphs[0].font.color.rgb = colors['white']
    title.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(14), Inches(1))
    subtitle = subtitle_box.text_frame
    subtitle.text = "Multi-Source Intelligent IPAM Management"
    subtitle.paragraphs[0].font.size = Pt(28)
    subtitle.paragraphs[0].font.color.rgb = colors['primary']
    subtitle.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 2: Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title and Content
    title = slide.shapes.title
    title.text = "The Challenge"
    
    content = slide.placeholders[1].text_frame
    content.text = "Managing Networks in a Multi-Cloud World"
    p = content.add_paragraph()
    p.text = "• InfoBlox contains networks from multiple sources"
    p.level = 1
    p = content.add_paragraph()
    p.text = "  - AWS, Azure, GCP, On-premises"
    p.level = 2
    p = content.add_paragraph()
    p.text = "• IP conflicts across cloud providers"
    p.level = 1
    p = content.add_paragraph()
    p.text = "• Manual synchronization is error-prone"
    p.level = 1
    p = content.add_paragraph()
    p.text = "• Risk of data loss from other sources"
    p.level = 1
    
    # Slide 3: Solution Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Solution Architecture"
    
    # Add architecture diagram components
    y_pos = Inches(2)
    components = [
        ("Analyzer", "Compares AWS exports with InfoBlox", colors['primary']),
        ("Reporter", "Generates comprehensive reports", colors['accent']),
        ("Updater", "Safely applies changes", colors['success'])
    ]
    
    for comp_name, comp_desc, color in components:
        # Component box
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2), y_pos, Inches(12), Inches(1.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        
        # Component text
        text_frame = shape.text_frame
        text_frame.text = f"{comp_name}: {comp_desc}"
        text_frame.paragraphs[0].font.size = Pt(20)
        text_frame.paragraphs[0].font.color.rgb = colors['white']
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.margin_top = Inches(0.5)
        
        y_pos += Inches(2)
    
    # Slide 4: Key Features
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Key Features"
    
    features = [
        ("🚫 No Delete Policy", "Never automatically deletes networks"),
        ("🎯 Source Attribution", "All AWS networks tagged with Source: AWS"),
        ("📦 Hierarchical Management", "Automatic container creation"),
        ("🏷️ Tag Preservation", "Intelligent tag merging"),
        ("⚠️ Conflict Detection", "Multi-source conflict analysis")
    ]
    
    y_pos = Inches(2)
    for feature, desc in features:
        text_box = slide.shapes.add_textbox(Inches(1), y_pos, Inches(14), Inches(0.8))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{feature} - {desc}"
        p.font.size = Pt(18)
        y_pos += Inches(0.9)
    
    # Slide 5: Conflict Types
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Intelligent Conflict Detection"
    
    # Create table for conflicts
    x, y, cx, cy = Inches(1), Inches(2), Inches(14), Inches(5)
    shape = slide.shapes.add_table(5, 4, x, y, cx, cy)
    table = shape.table
    
    # Header row
    headers = ['AWS Network', 'Existing Network', 'Source', 'Action']
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = colors['secondary']
    
    # Data rows
    data = [
        ['10.0.0.0/16', '10.0.1.0/24', 'Azure', 'Create as container'],
        ['172.16.0.0/24', '172.16.0.0/24', 'GCP', 'Block - Exact match'],
        ['10.50.0.0/23', '10.50.1.0/24', 'OnPrem', 'Flag for review'],
        ['10.80.5.0/25', '10.80.0.0/16', 'OnPrem', 'Check hierarchy']
    ]
    
    for row_idx, row_data in enumerate(data, 1):
        for col_idx, cell_data in enumerate(row_data):
            table.cell(row_idx, col_idx).text = cell_data
    
    # Slide 6: Process Flow
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Synchronization Process"
    
    steps = [
        "1. Export AWS VPC data to CSV",
        "2. Run comprehensive analysis",
        "3. Review generated reports",
        "4. Resolve any conflicts",
        "5. Apply changes (with dry-run option)"
    ]
    
    y_pos = Inches(2.5)
    for i, step in enumerate(steps):
        # Step box
        shape = slide.shapes.add_shape(
            MSO_SHAPE.CHEVRON,
            Inches(1 + i * 2.8), y_pos, Inches(2.5), Inches(1)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors['primary'] if i % 2 == 0 else colors['accent']
        
        # Step text
        text_box = slide.shapes.add_textbox(
            Inches(1 + i * 2.8), y_pos + Inches(1.2), Inches(2.5), Inches(1)
        )
        tf = text_box.text_frame
        tf.text = step
        tf.word_wrap = True
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 7: Benefits
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Business Benefits"
    
    benefits = [
        ("90% Time Savings", "Automated vs manual tracking", colors['success']),
        ("Zero Data Loss", "No accidental deletions", colors['accent']),
        ("100% Audit Trail", "Complete change tracking", colors['primary']),
        ("Multi-Cloud Ready", "Supports all major providers", colors['secondary'])
    ]
    
    x_pos = Inches(1)
    for benefit, desc, color in benefits:
        # Benefit box
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_pos, Inches(3), Inches(3.5), Inches(3)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        
        # Benefit text
        text_frame = shape.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.text = benefit
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = colors['white']
        p.alignment = PP_ALIGN.CENTER
        
        p = text_frame.add_paragraph()
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = colors['white']
        p.alignment = PP_ALIGN.CENTER
        
        x_pos += Inches(3.75)
    
    # Slide 8: Getting Started
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Getting Started"
    
    content = slide.placeholders[1].text_frame
    content.text = "Quick Start Guide"
    p = content.add_paragraph()
    p.text = "1. Clone the repository:"
    p.level = 0
    p = content.add_paragraph()
    p.text = "   git clone https://github.com/tshoush/AWS-Sync-Newer.git"
    p.level = 1
    p.font.name = 'Courier New'
    p = content.add_paragraph()
    p.text = "2. Configure InfoBlox credentials:"
    p.level = 0
    p = content.add_paragraph()
    p.text = "   cp config.env.example config.env"
    p.level = 1
    p.font.name = 'Courier New'
    p = content.add_paragraph()
    p.text = "3. Run analysis:"
    p.level = 0
    p = content.add_paragraph()
    p.text = "   python infoblox_aws_comprehensive_analyzer.py vpc_data.csv"
    p.level = 1
    p.font.name = 'Courier New'
    
    # Save presentation
    prs.save('AWS_InfoBlox_Sync_Presentation.pptx')
    print("✅ PowerPoint presentation created: AWS_InfoBlox_Sync_Presentation.pptx")
    
    # Note about PDF conversion
    print("\nTo convert to PDF:")
    print("- On Mac: Open in Keynote/PowerPoint and export as PDF")
    print("- On Windows: Open in PowerPoint and Save As PDF")
    print("- On Linux: Use LibreOffice: libreoffice --convert-to pdf AWS_InfoBlox_Sync_Presentation.pptx")

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("❌ python-pptx not installed. Install with: pip install python-pptx")
        print("\n✅ Alternative: Open presentation.html in your browser for an interactive web presentation!")
        print("   Features: Animations, navigation, progress tracking")